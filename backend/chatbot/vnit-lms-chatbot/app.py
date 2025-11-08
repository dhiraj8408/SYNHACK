import os
import random
import threading
import requests
import fitz  # PyMuPDF
import io
import tempfile
import mimetypes
import hashlib
from flask import Flask, request, jsonify
from flask_cors import CORS
import chromadb
from sentence_transformers import SentenceTransformer
from groq import Groq
from dotenv import load_dotenv
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Additional imports for multi-format support
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    print("Warning: python-pptx not installed. PPT support disabled.")

try:
    import pytesseract
    from PIL import Image
except ImportError:
    pytesseract = None
    Image = None
    print("Warning: pytesseract/Pillow not installed. Image OCR support disabled.")

try:
    from moviepy.editor import VideoFileClip
    import speech_recognition as sr
except ImportError:
    VideoFileClip = None
    sr = None
    print("Warning: Video processing libraries not installed. Video support disabled.")

# --- 1. LOAD ENVIRONMENT VARIABLES ---
print("Loading environment variables...")
load_dotenv()  # This loads the .env file

# --- 2. GET, SPLIT, AND CHOOSE API KEY ---
keys_string = os.environ.get("GROQ_API_KEYS")
if not keys_string:
    print("Error: GROQ_API_KEYS not found in .env file.")
    print("Please create a .env file and add: GROQ_API_KEYS='key1,key2,...'")
    exit()

api_key_list = keys_string.split(',')
selected_api_key = random.choice(api_key_list)
print(f"Using one of {len(api_key_list)} Groq API keys.")

# --- 3. SET UP FLASK APP AND CORS ---
app = Flask(__name__)
# This allows your teammates' frontend (on a different domain) to call your API
# Allow CORS from all origins (useful for development). If you need to restrict origins,
# replace '*' with a list of allowed origins or a specific domain.
CORS(app, resources={r"/*": {"origins": "*"}}, supports_credentials=True)

# --- 4. GLOBAL CONSTANTS ---
DB_PATH = "./chroma_db"
COLLECTION_NAME = "vnit_lms"
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
GROQ_MODEL = "llama-3.1-8b-instant" # Replaced the decommissioned model
EMBEDDING_MODEL = 'all-MiniLM-L6-v2'

# --- 5. INITIALIZE ALL MODELS AND DB (Load them once on startup) ---
print("--- Initializing models and connecting to DB... ---")
try:
    groq_client = Groq(api_key=selected_api_key)
    print("Groq client initialized.")
    
    embedding_model = SentenceTransformer(EMBEDDING_MODEL)
    print(f"Embedding model '{EMBEDDING_MODEL}' loaded.")
    
    db_client = chromadb.PersistentClient(path=DB_PATH)
    collection = db_client.get_or_create_collection(name=COLLECTION_NAME)
    print(f"Connected to ChromaDB at '{DB_PATH}'. Collection '{COLLECTION_NAME}' loaded.")
    
    print("--- Server is ready to receive requests. ---")
except Exception as e:
    print(f"--- FATAL STARTUP ERROR: {e} ---")
    print("Please check your API keys, model names, and file permissions.")
    exit()


# --- 6. INGESTION HELPER FUNCTIONS ---

def get_google_drive_file_id(url):
    """Extracts the file ID from various Google Drive share link formats."""
    try:
        # Format: /file/d/FILE_ID/view
        if '/d/' in url:
            return url.split('/d/')[1].split('/')[0]
        # Format: ?id=FILE_ID
        if 'id=' in url:
            return url.split('id=')[1].split('&')[0]
        print(f"Could not parse File ID from URL: {url}")
        return None
    except Exception as e:
        print(f"Error parsing Google Drive URL: {e}")
        return None

def detect_file_type(file_content, filename=None, content_type=None):
    """Detects file type from content, filename, or content-type header."""
    # First, try to detect from content_type
    if content_type:
        if 'pdf' in content_type.lower():
            return 'pdf'
        elif 'powerpoint' in content_type.lower() or 'presentation' in content_type.lower():
            return 'ppt'
        elif 'image' in content_type.lower():
            return 'image'
        elif 'video' in content_type.lower():
            return 'video'
    
    # Try to detect from filename extension
    if filename:
        ext = filename.lower().split('.')[-1] if '.' in filename else ''
        if ext in ['pdf']:
            return 'pdf'
        elif ext in ['ppt', 'pptx']:
            return 'ppt'
        elif ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp']:
            return 'image'
        elif ext in ['mp4', 'avi', 'mov', 'mkv', 'wmv', 'flv', 'webm']:
            return 'video'
    
    # Try to detect from file content (magic bytes)
    if file_content:
        # PDF magic bytes: %PDF
        if file_content[:4] == b'%PDF':
            return 'pdf'
        # PPTX is a ZIP file (starts with PK)
        elif file_content[:2] == b'PK':
            # Check if it's a PPTX by looking for specific structure
            if b'ppt/' in file_content[:1024] or b'[Content_Types].xml' in file_content[:1024]:
                return 'ppt'
        # Image formats
        elif file_content[:2] == b'\xff\xd8':  # JPEG
            return 'image'
        elif file_content[:8] == b'\x89PNG\r\n\x1a\n':  # PNG
            return 'image'
        elif file_content[:6] in [b'GIF87a', b'GIF89a']:  # GIF
            return 'image'
        # Video formats (basic detection)
        elif file_content[:4] in [b'ftyp', b'\x00\x00\x00\x20']:  # MP4
            return 'video'
        elif file_content[:4] == b'RIFF':  # AVI
            return 'video'
    
    return 'unknown'

def extract_text_from_pdf(file_content):
    """Extracts text from PDF file content using OCR (handles handwritten PDFs)."""
    if pytesseract is None or Image is None:
        print("Error: pytesseract/Pillow libraries not available for PDF OCR.")
        return None
    
    try:
        full_text = ""
        with io.BytesIO(file_content) as pdf_stream:
            with fitz.open(stream=pdf_stream, filetype="pdf") as doc:
                print(f"[PDF Processing] Processing {len(doc)} pages with OCR...")
                for page_num, page in enumerate(doc):
                    print(f"[PDF Processing] OCR on page {page_num + 1}/{len(doc)}...")
                    # Convert PDF page to image (pixmap)
                    # Use a high DPI for better OCR accuracy
                    mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better quality
                    pix = page.get_pixmap(matrix=mat)
                    
                    # Convert pixmap to PIL Image
                    img_data = pix.tobytes("png")
                    page_image = Image.open(io.BytesIO(img_data))
                    
                    # Convert to RGB if necessary
                    if page_image.mode != 'RGB':
                        page_image = page_image.convert('RGB')
                    
                    # Perform OCR on the page
                    page_text = pytesseract.image_to_string(page_image)
                    if page_text.strip():
                        full_text += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
        
        return full_text if full_text.strip() else None
    except Exception as e:
        print(f"Error extracting text from PDF with OCR: {e}")
        import traceback
        traceback.print_exc()
        return None

def extract_text_from_ppt(file_content):
    """Extracts text from PowerPoint file content."""
    if Presentation is None:
        print("Error: python-pptx library not available.")
        return None
    
    try:
        full_text = ""
        with io.BytesIO(file_content) as ppt_stream:
            prs = Presentation(ppt_stream)
            for slide_num, slide in enumerate(prs.slides):
                full_text += f"\n--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text + "\n"
                    # Also check for tables
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            row_text = " | ".join([cell.text for cell in row.cells])
                            full_text += row_text + "\n"
        return full_text
    except Exception as e:
        print(f"Error extracting text from PPT: {e}")
        return None

def extract_text_from_image(file_content):
    """Extracts text from image using OCR."""
    if pytesseract is None or Image is None:
        print("Error: pytesseract/Pillow libraries not available.")
        return None
    
    try:
        image = Image.open(io.BytesIO(file_content))
        # Convert to RGB if necessary
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        # Perform OCR
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        print(f"Error extracting text from image: {e}")
        return None

def extract_text_from_video(file_content):
    """Extracts text from video by transcribing audio only."""
    if VideoFileClip is None or sr is None:
        print("Error: Video processing libraries not available.")
        return None
    
    try:
        full_text = ""
        
        # Save video to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.mp4') as tmp_video:
            tmp_video.write(file_content)
            tmp_video_path = tmp_video.name
        
        try:
            # Extract and transcribe audio
            print("[Video Processing] Extracting audio...")
            video = VideoFileClip(tmp_video_path)
            
            if video.audio is None:
                print("[Video Processing] Warning: Video has no audio track.")
                video.close()
                return None
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.wav') as tmp_audio:
                tmp_audio_path = tmp_audio.name
                video.audio.write_audiofile(tmp_audio_path, verbose=False, logger=None)
            
            video.close()
            
            # Transcribe audio
            print("[Video Processing] Transcribing audio...")
            recognizer = sr.Recognizer()
            with sr.AudioFile(tmp_audio_path) as source:
                # Adjust for ambient noise
                recognizer.adjust_for_ambient_noise(source)
                audio_data = recognizer.record(source)
                
                try:
                    transcript = recognizer.recognize_google(audio_data)
                    full_text = transcript
                    print(f"[Video Processing] Transcription completed ({len(transcript)} characters).")
                except sr.UnknownValueError:
                    print("[Video Processing] Could not understand audio (audio may be too quiet or unclear).")
                    full_text = None
                except sr.RequestError as e:
                    print(f"[Video Processing] Error with speech recognition service: {e}")
                    full_text = None
            
            # Clean up audio file
            if os.path.exists(tmp_audio_path):
                os.remove(tmp_audio_path)
            
        finally:
            # Clean up video file
            if os.path.exists(tmp_video_path):
                os.remove(tmp_video_path)
        
        return full_text if full_text and full_text.strip() else None
        
    except Exception as e:
        print(f"Error extracting text from video: {e}")
        import traceback
        traceback.print_exc()
        return None

def process_file(file_content, file_type, file_id=None, filename=None):
    """
    Unified background task to process any file type:
    1. Detects file type (if not provided)
    2. Extracts text based on file type
    3. Chunks the text
    4. Embeds and stores in ChromaDB
    """
    print(f"[Background Ingest] Starting processing. File type: {file_type}, Size: {len(file_content)} bytes")
    
    try:
        # 1. Detect file type if not provided
        if file_type == 'unknown' or file_type is None:
            file_type = detect_file_type(file_content, filename=filename)
            print(f"[Background Ingest] Detected file type: {file_type}")
        
        if file_type == 'unknown':
            print("[Background Ingest] Error: Could not detect file type.")
            return
        
        # 2. Extract text based on file type
        print(f"[Background Ingest] Extracting text from {file_type}...")
        full_text = None
        
        if file_type == 'pdf':
            full_text = extract_text_from_pdf(file_content)
        elif file_type == 'ppt':
            full_text = extract_text_from_ppt(file_content)
        elif file_type == 'image':
            full_text = extract_text_from_image(file_content)
        elif file_type == 'video':
            full_text = extract_text_from_video(file_content)
        else:
            print(f"[Background Ingest] Error: Unsupported file type: {file_type}")
            return
        
        if not full_text or not full_text.strip():
            print(f"[Background Ingest] Error: No text extracted from {file_type} file.")
            return
        
        print(f"[Background Ingest] Extracted {len(full_text)} characters of text.")
        
        # 3. Chunk the text
        print("[Background Ingest] Chunking text...")
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            length_function=len
        )
        text_chunks = text_splitter.split_text(full_text)
        print(f"[Background Ingest] Split into {len(text_chunks)} chunks.")
        
        if not text_chunks:
            print("[Background Ingest] Error: Text could not be split into chunks.")
            return
        
        # 4. Embed and Store
        print(f"[Background Ingest] Generating embeddings for {len(text_chunks)} chunks...")
        embeddings = embedding_model.encode(text_chunks, show_progress_bar=True)
        
        # Create unique IDs
        if file_id:
            ids = [f"{file_type}_{file_id}_chunk_{i}" for i in range(len(text_chunks))]
        else:
            file_hash = hashlib.md5(file_content).hexdigest()[:8]
            ids = [f"{file_type}_{file_hash}_chunk_{i}" for i in range(len(text_chunks))]
        
        print(f"[Background Ingest] Adding {len(ids)} chunks to ChromaDB...")
        collection.add(
            embeddings=embeddings.tolist(),
            documents=text_chunks,
            ids=ids
        )
        
        print(f"--- [Background Ingest] COMPLETED: Ingestion for {file_type} file ---")
        
    except Exception as e:
        print(f"--- [Background Ingest] FAILED: {e} ---")
        import traceback
        traceback.print_exc()

def process_drive_link(drive_url):
    """
    Downloads file from Google Drive and processes it.
    Supports PDF, PPT, images, and videos.
    """
    print(f"[Background Ingest] Starting task for Google Drive link: {drive_url}")
    
    try:
        # 1. Get File ID and create download link
        file_id = get_google_drive_file_id(drive_url)
        if not file_id:
            print("[Background Ingest] Error: Could not parse Google Drive File ID.")
            return
        
        download_url = f"https://drive.google.com/uc?export=download&id={file_id}"
        
        # 2. Download the file
        print(f"[Background Ingest] Downloading file: {file_id}...")
        response = requests.get(download_url)
        
        # Handle Google's large file confirmation redirect
        if "confirm=t" in response.url:
            print("[Background Ingest] Large file. Following confirmation redirect...")
            response = requests.get(download_url, cookies=response.cookies)
        
        # Check if the download failed
        if response.status_code != 200:
            print(f"[Background Ingest] Error: Failed to download file. Status: {response.status_code}")
            print("Please ensure the link is public ('Anyone with the link').")
            return
        
        # Check if we actually got a file
        if not response.content:
            print("[Background Ingest] Error: Downloaded file is empty.")
            return
        
        content_type = response.headers.get('content-type', '')
        print(f"[Background Ingest] File downloaded ({len(response.content)} bytes). Content-Type: {content_type}")
        
        # 3. Detect file type and process
        file_type = detect_file_type(response.content, content_type=content_type)
        process_file(response.content, file_type, file_id=file_id)
        
    except Exception as e:
        print(f"--- [Background Ingest] FAILED: {e} ---")
        import traceback
        traceback.print_exc()


# --- 7. API ENDPOINT: /chatbot-api/ingest (For Professors) ---

@app.route("/chatbot-api/ingest", methods=["POST"])
def handle_ingestion():
    """
    API endpoint to trigger the ingestion of a new document from Google Drive.
    Accepts Google Drive links (via 'drive_link' in JSON).
    Supports: PDF, PPT/PPTX, Images, and Videos.
    
    Returns an "Accepted" response immediately and processes in background.
    """
    print(f"\n--- New Ingestion Request Received ---")
    
    # Get JSON data
    data = request.json
    if not data:
        return jsonify({"error": "No JSON data provided"}), 400
    
    drive_link = data.get("drive_link")
    
    if not drive_link:
        return jsonify({"error": "No 'drive_link' provided. Please provide a Google Drive link in JSON format: {\"drive_link\": \"...\"}"}), 400
    
    print(f"Google Drive link: {drive_link}")
    
    # Start the ingestion process in a separate thread
    # This allows the API to return an immediate response
    thread = threading.Thread(target=process_drive_link, args=(drive_link,))
    thread.start()
    
    # Return a 202 "Accepted" status
    return jsonify({
        "message": "Ingestion process started. The document will be available for queries in a few minutes."
    }), 202


# --- 8. API ENDPOINT: /chatbot-api/chat (For Students) ---

@app.route("/chatbot-api/chat", methods=["POST"])
def handle_chat():
    """
    Handles chat requests from the frontend.
    Performs RAG to answer questions based on ingested documents.
    """
    
    data = request.json
    user_question = data.get("question")

    if not user_question:
        return jsonify({"error": "No question provided"}), 400

    print(f"\n--- New Chat Request ---")
    print(f"Question: {user_question}")

    try:
        # 1. Vectorize the user's question
        question_vector = embedding_model.encode(user_question).tolist()

        # 2. Query ChromaDB to find relevant context
        print("Searching for context...")
        results = collection.query(
            query_embeddings=[question_vector],
            n_results=5 # Get top 5 most relevant chunks
        )
        
        context_chunks = results['documents'][0]
        context = "\n\n".join(context_chunks)
        
        if not context_chunks:
            print("No relevant context found in database.")
            return jsonify({"answer": "I'm sorry, but I don't have that information in my knowledge base."})

        print(f"Found context: {context[:200]}...") # Print first 200 chars

        # 3. Build the prompt for Groq
        system_prompt = """
        You are a helpful AI teaching assistant for the VNIT Learning Management System.
        Your goal is to answer student questions based ONLY on the provided course materials.
        
        Follow these rules STRICTLY:
        1. Use ONLY the 'Provided Context' below to answer the question.
        2. If the answer is not in the context, clearly state: "I'm sorry, but I don't have that information in my knowledge base."
        3. Do not make up answers or use any external knowledge.
        4. Be concise and direct in your response.
        """
        
        final_prompt = f"""
        {system_prompt}
        
        ---
        Provided Context:
        {context}
        ---
        
        Student Question:
        {user_question}
        
        Answer:
        """

        # 4. Call Groq LLM
        print("Sending prompt to Groq...")
        chat_completion = groq_client.chat.completions.create(
            messages=[
                {"role": "user", "content": final_prompt}
            ],
            model=GROQ_MODEL,
        )
        
        bot_answer = chat_completion.choices[0].message.content
        print(f"Groq Answer: {bot_answer}")

        # 5. Send Response
        return jsonify({"answer": bot_answer})

    except Exception as e:
        print(f"Error during RAG pipeline: {e}")
        return jsonify({"error": f"An internal error occurred: {e}"}), 500


# --- 9. RUN THE FLASK SERVER ---
if __name__ == "__main__":
    # Runs the server on port 5000 and makes it accessible on your network
    # 'debug=True' automatically reloads the server when you save changes
    app.run(host='0.0.0.0', port=5001, debug=True)